import glob
import os
from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
SLIDES_DIR = os.path.join(HERE, "slides")
OUT = os.path.join(os.path.dirname(HERE), "Fakerni_Defense_Deck.pptx")

# 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

pngs = sorted(glob.glob(os.path.join(SLIDES_DIR, "slide_*.png")))
if not pngs:
    raise SystemExit("no slide PNGs found")

for png in pngs:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)

prs.save(OUT)
print(f"wrote {OUT} with {len(pngs)} slides")
